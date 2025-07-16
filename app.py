from flask import Flask, request, render_template, send_file, jsonify
from werkzeug.utils import secure_filename
from pptx import Presentation
from googletrans import Translator
import os
import tempfile
import time
import threading
import uuid

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB 제한
app.config['SECRET_KEY'] = 'your-secret-key-here'
translator = Translator()

# 진행률 추적을 위한 전역 변수
progress_data = {}
# 중단 신호를 위한 전역 변수
cancel_signals = {}

# 지원하는 언어 목록 (PPTX 번역 지원 언어)
SUPPORTED_LANGUAGES = {
    'ko': '한국어',
    'en': 'English',
    'ja': '日本語',
    'zh-cn': '中文(简体)',
    'zh-tw': '中文(繁體)',
    'es': 'Español',
    'fr': 'Français',
    'de': 'Deutsch',
    'ru': 'Русский',
    'pt': 'Português',
    'it': 'Italiano'
}

# 다국어 텍스트
TEXTS = {
    'ko': {
        'title': 'PPTX 번역기 - PPTX 파일 번역 서비스',
        'subtitle': 'PPTX 파일의 일본어 텍스트를 다양한 언어로 번역해드립니다',
        'upload_title': 'PPTX 파일을 선택하거나 드래그하세요',
        'upload_subtitle': '지원 형식: .pptx (최대 16MB)',
        'source_lang': '원본 언어',
        'target_lang': '번역할 언어',
        'auto_detect': '자동 감지',
        'auto_detect_note': '일본어만 지원 (자동 감지)',
        'select_lang': '언어를 선택하세요',
        'translate_btn': '번역 시작',
        'translating': '번역 중...',
        'progress_title': '번역 중입니다...',
        'progress_subtitle': '파일 크기와 내용에 따라 시간이 걸릴 수 있습니다',
        'feature_excel': 'PPTX 지원',
        'feature_excel_desc': '.pptx 파일 완벽 지원',
        'feature_multilang': '다국어 번역',
        'feature_multilang_desc': '10개 언어 간 자유로운 번역',
        'feature_safe': '안전한 처리',
        'feature_safe_desc': '임시 파일 자동 삭제',
        'language_switch': '한국어',
        'progress_cell': '텍스트',
        'progress_sheet': '슬라이드',
        'progress_eta': '예상 종료 시간',
        'progress_elapsed': '경과 시간',
        'progress_remaining': '남은 시간',
        'file_selected': '파일이 선택되었습니다',
        'file_remove': '파일 해제',
        'file_remove_confirm': '선택된 파일을 해제하시겠습니까?',
        'file_required': '파일을 선택해주세요.',
        'target_lang_required': '번역할 언어를 선택해주세요.',
        'translation_complete': '번역이 완료되었습니다! 파일이 다운로드됩니다.',
        'translation_error': '번역 중 오류가 발생했습니다.',
        'file_error': '오류',
        'file_not_found': '파일을 찾을 수 없습니다.',
        'invalid_pptx': '유효하지 않은 PPTX 파일입니다',
        'file_open_error': '파일 열기 실패',
        'pptx_only': 'PPTX 파일(.pptx)만 지원합니다.',
        'unsupported_lang': '지원하지 않는 언어입니다.',
        'translation_failed': '번역 실패',
        'download_ready': '다운로드 준비 완료',
        'sheet_name_translation': '슬라이드 제목 번역 포함',
        'sheet_name_translation_desc': '슬라이드 제목도 함께 번역됩니다',
        'cancel_translation': '번역 중단',
        'cancel_confirm': '번역을 중단하시겠습니까?',
        'translation_cancelled': '번역이 중단되었습니다.',
        'cancelling': '중단 중...'
    },
    'en': {
        'title': 'PPTX Translator - PPTX File Translation Service',
        'subtitle': 'Translate Japanese text in PPTX files into various languages',
        'upload_title': 'Select or drag PPTX file here',
        'upload_subtitle': 'Supported format: .pptx (max 16MB)',
        'source_lang': 'Source Language',
        'target_lang': 'Target Language',
        'auto_detect': 'Auto Detect',
        'auto_detect_note': 'Only Japanese supported (auto detect)',
        'select_lang': 'Select language',
        'translate_btn': 'Start Translation',
        'translating': 'Translating...',
        'progress_title': 'Translation in progress...',
        'progress_subtitle': 'Time may vary depending on file size and content',
        'feature_excel': 'PPTX Support',
        'feature_excel_desc': 'Perfect support for .pptx files',
        'feature_multilang': 'Multi-language Translation',
        'feature_multilang_desc': 'Free translation between 10 languages',
        'feature_safe': 'Safe Processing',
        'feature_safe_desc': 'Automatic temporary file deletion',
        'language_switch': 'English',
        'progress_cell': 'text',
        'progress_sheet': 'slide',
        'progress_eta': 'Estimated completion time',
        'progress_elapsed': 'Elapsed time',
        'progress_remaining': 'Remaining time',
        'file_selected': 'File selected',
        'file_remove': 'Remove File',
        'file_remove_confirm': 'Are you sure you want to remove the selected file?',
        'file_required': 'Please select a file.',
        'target_lang_required': 'Please select a target language.',
        'translation_complete': 'Translation completed! File will be downloaded.',
        'translation_error': 'An error occurred during translation.',
        'file_error': 'Error',
        'file_not_found': 'File not found.',
        'invalid_pptx': 'Invalid PPTX file',
        'file_open_error': 'File open failed',
        'pptx_only': 'Only PPTX files (.pptx) are supported.',
        'unsupported_lang': 'Unsupported language.',
        'translation_failed': 'Translation failed',
        'download_ready': 'Download ready',
        'sheet_name_translation': 'Slide Title Translation',
        'sheet_name_translation_desc': 'Slide titles are also translated',
        'cancel_translation': 'Cancel Translation',
        'cancel_confirm': 'Are you sure you want to cancel the translation?',
        'translation_cancelled': 'Translation has been cancelled.',
        'cancelling': 'Cancelling...'
    }
}

def is_japanese(text):
    return any('\u3040' <= ch <= '\u30ff' or '\u4e00' <= ch <= '\u9faf' for ch in text)

def translate_pptx_file_with_progress(file_path, target_lang, task_id=None):
    translator = Translator()
    prs = Presentation(file_path)
    total_texts = 0
    translated_texts = 0
    current_text = 0
    slides = list(prs.slides)
    # 전체 텍스트 수 계산
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if is_japanese(run.text):
                            total_texts += 1
    start_time = time.time()
    if task_id:
        progress_data[task_id] = {
            'progress': 0,
            'current_cell': 0,
            'total_cells': total_texts,
            'translated_cells': 0,
            'current_sheet': 1,
            'sheet_progress': f"0/{len(slides)}",
            'elapsed_time': 0,
            'remaining_time': 0,
            'eta': '00:00:00',
            'status': 'running'
        }
        cancel_signals[task_id] = False
    try:
        for slide_idx, slide in enumerate(slides):
            if task_id and cancel_signals.get(task_id, False):
                if task_id:
                    progress_data[task_id]['status'] = 'cancelled'
                return None, "Translation cancelled by user"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if is_japanese(run.text):
                                current_text += 1
                                try:
                                    run.text = translator.translate(run.text, src='ja', dest=target_lang).text
                                    translated_texts += 1
                                except Exception as e:
                                    print(f"번역 실패: {run.text} ({e})")
                                # 진행률 업데이트
                                if task_id and (current_text % max(1, total_texts // 10) == 0 or current_text == total_texts):
                                    progress = (current_text / total_texts) * 100
                                    elapsed_time = time.time() - start_time
                                    if current_text > 0:
                                        estimated_total_time = (elapsed_time / current_text) * total_texts
                                        remaining_time = estimated_total_time - elapsed_time
                                    else:
                                        remaining_time = 0
                                    progress_data[task_id].update({
                                        'progress': round(progress, 1),
                                        'current_cell': current_text,
                                        'translated_cells': translated_texts,
                                        'current_sheet': slide_idx + 1,
                                        'sheet_progress': f"{slide_idx + 1}/{len(slides)}",
                                        'elapsed_time': round(elapsed_time, 1),
                                        'remaining_time': round(remaining_time, 1),
                                        'eta': time.strftime('%H:%M:%S', time.gmtime(remaining_time)) if remaining_time > 0 else '00:00:00'
                                    })
        if task_id:
            progress_data[task_id]['status'] = 'completed'
            progress_data[task_id]['progress'] = 100
        output_path = tempfile.mktemp(suffix="_translated.pptx")
        prs.save(output_path)
        return output_path, f"{translated_texts} texts translated"
    except Exception as e:
        return None, f"번역 실패: {e}"

def translate_pptx_file(file_path, target_lang):
    translator = Translator()
    prs = Presentation(file_path)
    translated_texts = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if is_japanese(run.text):
                            try:
                                run.text = translator.translate(run.text, src='ja', dest=target_lang).text
                                translated_texts += 1
                            except Exception as e:
                                print(f"번역 실패: {run.text} ({e})")
    output_path = tempfile.mktemp(suffix="_translated.pptx")
    prs.save(output_path)
    return output_path, f"{translated_texts} texts translated"

@app.route('/')
def index():
    lang = request.args.get('lang', 'en')
    if lang not in TEXTS:
        lang = 'en'
    return render_template('index.html', languages=SUPPORTED_LANGUAGES, texts=TEXTS[lang], current_lang=lang)

@app.route('/translate', methods=['POST'])
def translate_file():
    lang = request.form.get('lang', 'en')
    if lang not in TEXTS:
        lang = 'en'
    if 'file' not in request.files:
        return jsonify({'error': TEXTS[lang]['file_required']}), 400
    file = request.files['file']
    target_lang = request.form.get('target_lang', 'ko')
    if file.filename == '':
        return jsonify({'error': TEXTS[lang]['file_required']}), 400
    if not file.filename.endswith('.pptx'):
        return jsonify({'error': TEXTS[lang]['pptx_only']}), 400
    if target_lang not in SUPPORTED_LANGUAGES:
        return jsonify({'error': TEXTS[lang]['unsupported_lang']}), 400
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
            file.save(tmp_file.name)
            tmp_path = tmp_file.name
        task_id = request.form.get('task_id', None)
        output_path = None
        try:
            if task_id:
                output_path, message = translate_pptx_file_with_progress(tmp_path, target_lang, task_id)
            else:
                output_path, message = translate_pptx_file(tmp_path, target_lang)
            if output_path is None:
                try:
                    os.unlink(tmp_path)
                except:
                    pass
                if "Translation cancelled by user" in message:
                    error_msg = TEXTS[lang]['translation_cancelled']
                elif "유효하지 않은 PPTX 파일" in message:
                    error_msg = TEXTS[lang]['invalid_pptx']
                elif "파일 열기 실패" in message:
                    error_msg = TEXTS[lang]['file_open_error']
                else:
                    error_msg = message
                return jsonify({'error': error_msg}), 400
        except Exception as e:
            try:
                os.unlink(tmp_path)
            except:
                pass
            raise e
        try:
            os.unlink(tmp_path)
        except Exception as e:
            print(f"Warning: Could not delete temporary file {tmp_path}: {e}")
        return send_file(
            output_path,
            as_attachment=True,
            download_name="translated_" + secure_filename(file.filename),
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        return jsonify({'error': f"{TEXTS[lang]['translation_error']}: {str(e)}"}), 500

@app.route('/progress/<task_id>')
def get_progress(task_id):
    if task_id in progress_data:
        return jsonify(progress_data[task_id])
    else:
        return jsonify({'error': 'Task not found'}), 404

@app.route('/start_translation', methods=['POST'])
def start_translation():
    task_id = str(uuid.uuid4())
    progress_data[task_id] = {
        'progress': 0,
        'status': 'starting'
    }
    cancel_signals[task_id] = False
    return jsonify({'task_id': task_id})

@app.route('/cancel_translation/<task_id>', methods=['POST'])
def cancel_translation(task_id):
    if task_id in cancel_signals:
        cancel_signals[task_id] = True
        return jsonify({'status': 'cancellation_requested'})
    else:
        return jsonify({'error': 'Task not found'}), 404

@app.route('/health')
def health_check():
    return jsonify({'status': 'healthy'})

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=int(os.environ.get('PORT', 5000))) 